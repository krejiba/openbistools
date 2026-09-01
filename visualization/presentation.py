from pptx import Presentation
from docx import Document
import pdfplumber
from PIL import Image
import io


def convert_doc_to_images(doc_path):
    doc = Document(doc_path)
    pil_images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img_bytes = rel.target_part.blob
            try:
                pil_img = Image.open(io.BytesIO(img_bytes))
                pil_images.append(pil_img.copy())
                pil_img.close()
            except Exception as e:
                print(f"Could not open image: {e}")

    return pil_images


def convert_ppt_to_images(ppt_file):
    prs = Presentation(ppt_file)
    image_files = []
    for slide_number, slide in enumerate(prs.slides):
        img = Image.new(
            "RGB", (int(prs.slide_width.pt), int(prs.slide_height.pt)), "white"
        )
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                img_stream = io.BytesIO(shape.image.blob)
                shape_img = Image.open(img_stream)
                left = int(shape.left.pt)
                top = int(shape.top.pt)
                img.paste(shape_img, (left, top))
        image_files.append(img)
    return image_files


def convert_pdf_to_images(pdf_file, resolution=150):
    images = []
    with pdfplumber.open(pdf_file) as pdf:
        for page in pdf.pages:
            page_image = page.to_image(resolution=resolution)
            pil_img = page_image.original.convert("RGB")
            images.append(pil_img)
    return images


def create_image_grid(images, grid_size=(3, 3)):
    if not images:
        return
    images = images[: grid_size[0] * grid_size[1]]
    width, height = images[0].size
    grid_width = width * grid_size[0]
    grid_height = height * grid_size[1]
    new_image = Image.new("RGB", (grid_width, grid_height), "white")
    for index, img in enumerate(images):
        x_offset = (index % grid_size[0]) * width
        y_offset = (index // grid_size[0]) * height
        new_image.paste(img, (x_offset, y_offset))
    return new_image


def create_grid_from_ppt(ppt_file, grid_size=(3, 3)):
    images = convert_ppt_to_images(ppt_file)
    return create_image_grid(images, grid_size)


def create_grid_from_doc(doc_file, grid_size=(3, 3)):
    images = convert_doc_to_images(doc_file)
    return create_image_grid(images, grid_size)


def create_grid_from_pdf(pdf_file, grid_size=(3, 3), resolution=150):
    images = convert_pdf_to_images(pdf_file, resolution=resolution)
    return create_image_grid(images, grid_size)
